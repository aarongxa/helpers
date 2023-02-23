#!/usr/bin/python

import sys
import codecs
from pptx import Presentation

# Check that both arguments were provided
if len(sys.argv) != 3:
    print("Usage: python pptx_to_txt.py input.pptx output.txt")
    sys.exit(1)

# Open the PowerPoint file
pptx_file = sys.argv[1]
prs = Presentation(pptx_file)

# Open the output text file in Unicode mode
txt_file = sys.argv[2]
with codecs.open(txt_file, "w", encoding="utf-8") as output:
    # Loop through each slide in the presentation
    for slide in prs.slides:
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape is a text box
            if shape.has_text_frame:
                # Write the text to the output file
                output.write(shape.text.strip() + "\n")